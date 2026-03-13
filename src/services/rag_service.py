from __future__ import annotations

import io
import time
from pathlib import Path
from typing import Any

from google import genai
from google.genai import errors as genai_errors
from google.genai import types

from src.services.config_service import (
    GEMINI_API_KEY,
    SUMMARY_REQUIRED,
    UPLOAD_TIMEOUT_SECONDS,
)
from src.services.logging_service import get_logger

SUMMARY_PROMPT = """You are a retrieval optimization assistant.
Analyze the document and generate a structured summary using EXACTLY this markdown schema.
Every section must appear in this exact order. Do not skip any section.

# Document Summary

## Document Title
[Clear human-readable title]

## Document Type
[Choose one: architecture | technical_spec | meeting_notes | policy | resume | tutorial | api_reference | research | general_document]

## Document Purpose
[2-3 sentences: why this document exists and what it explains or achieves]

## Main Topics
[Bullet list of primary subjects covered]

## Key Entities
### People
[Full names of any individuals referenced]
### Systems / Platforms
[Products, services, platforms, infrastructure]
### Tools / Frameworks
[Libraries, SDKs, APIs, tools referenced]

## Key Facts
[Bullet list of important facts: decisions, configurations, rules, metrics, limitations]

## Searchable Questions
[List 5 questions a user might ask that this document can answer. Write them as real user queries.]

## Short Abstract
[One dense paragraph: key concepts, systems mentioned, main purpose of the document]
"""

logger = get_logger(__name__)


class RagServiceError(Exception):
    """Base exception for RAG service failures."""


class RagUploadTimeoutError(RagServiceError):
    """Raised when Google indexing exceeds the shared upload budget."""


class RagSummaryError(RagServiceError):
    """Raised when summary generation fails and fallback is disabled."""


class RagService:
    def __init__(self, client: Any | None = None) -> None:
        self._client = client

    @property
    def client(self) -> Any:
        if self._client is None:
            self._client = genai.Client(api_key=GEMINI_API_KEY)
        return self._client

    def create_store(self, display_name: str) -> dict[str, str]:
        try:
            store = self.client.file_search_stores.create(
                config={"display_name": display_name}
            )
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        logger.info("created store display_name=%s store_id=%s", display_name, store.name)
        return {
            "name": store.name,
            "display_name": getattr(store, "display_name", display_name),
        }

    def delete_store(self, store_name: str) -> None:
        try:
            self.client.file_search_stores.delete(name=store_name)
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        logger.info("deleted store store_id=%s", store_name)

    def list_documents(self, store_name: str) -> list[dict[str, str]]:
        try:
            documents = list(self.client.file_search_stores.documents.list(parent=store_name))
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        return [
            {
                "name": doc.name,
                "display_name": getattr(doc, "display_name", ""),
            }
            for doc in documents
        ]

    def get_store_details(self, store_name: str) -> dict[str, Any]:
        try:
            store = self.client.file_search_stores.get(name=store_name)
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        return {
            "name": getattr(store, "name", store_name),
            "display_name": getattr(store, "display_name", ""),
            "raw": self._to_jsonable(store),
        }

    def verify_stores(self) -> list[dict[str, Any]]:
        try:
            stores = list(self.client.file_search_stores.list())
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        verified: list[dict[str, Any]] = []
        for store in stores:
            try:
                documents = list(self.client.file_search_stores.documents.list(parent=store.name))
            except genai_errors.APIError as exc:
                raise RagServiceError(str(exc)) from exc
            verified.append(
                {
                    "name": getattr(store, "name", ""),
                    "display_name": getattr(store, "display_name", ""),
                    "document_count": len(documents),
                }
            )
        return verified

    def get_document_details(self, document_name: str) -> dict[str, Any]:
        try:
            document = self.client.file_search_stores.documents.get(name=document_name)
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        return {
            "name": getattr(document, "name", document_name),
            "display_name": getattr(document, "display_name", ""),
            "raw": self._to_jsonable(document),
        }

    def get_operation_status(self, operation_name: str) -> dict[str, Any]:
        try:
            operation = types.UploadToFileSearchStoreOperation(name=operation_name)
            refreshed = self.client.operations.get(operation)
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        return {
            "name": getattr(refreshed, "name", operation_name),
            "done": getattr(refreshed, "done", None),
            "metadata": self._to_jsonable(getattr(refreshed, "metadata", None)),
            "error": self._to_jsonable(getattr(refreshed, "error", None)),
        }

    @staticmethod
    def _extract_text_for_summary(file_bytes: bytes, mime_type: str) -> tuple[bytes, str]:
        """Return (bytes, mime_type) safe for Gemini generate_content.

        PDF and text/* are passed through unchanged.
        Office formats (docx, xlsx, pptx) are converted to plain text.
        Extracted text is capped at 8000 characters to avoid token limits.
        """
        _DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        _XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        _PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        if mime_type == "application/pdf" or mime_type.startswith("text/"):
            return file_bytes, mime_type

        if mime_type == _DOCX:
            import docx  # python-docx
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif mime_type == _XLSX:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            rows: list[str] = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        rows.append(row_text)
            text = "\n".join(rows)
        elif mime_type == _PPTX:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slides: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slides.append(shape.text)
            text = "\n".join(slides)
        else:
            # application/json, application/rtf, etc. — decode as plain text
            text = file_bytes.decode("utf-8", errors="replace")

        text = text[:8000]
        return text.encode("utf-8"), "text/plain"

    def generate_summary(
        self,
        *,
        file_bytes: bytes,
        mime_type: str,
        prompt: str = SUMMARY_PROMPT,
        model: str = "gemini-2.5-flash",
    ) -> str:
        file_bytes, mime_type = self._extract_text_for_summary(file_bytes, mime_type)
        try:
            response = self.client.models.generate_content(
                model=model,
                contents=[
                    types.Part.from_bytes(data=file_bytes, mime_type=mime_type),
                    prompt,
                ],
            )
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        logger.info("generated summary mime_type=%s model=%s", mime_type, model)
        return response.text

    def upload_document_pair(
        self,
        *,
        store_name: str,
        original_path: str | Path,
        display_name: str,
        summary_text: str | None,
        summary_filename: str,
    ) -> dict[str, str]:
        if summary_text is None:
            if SUMMARY_REQUIRED:
                raise RagSummaryError(
                    "Summary generation failed and SUMMARY_REQUIRED=true."
                )
            logger.warning(
                "Summary generation unavailable; falling back to original-only upload."
            )

        start = time.perf_counter()
        original_doc = self._upload_and_wait(
            file_path=original_path,
            store_name=store_name,
            display_name=display_name,
            deadline=self._deadline_from_start(start),
        )

        summary_doc_name = ""
        if summary_text is not None:
            summary_path = self._write_summary_file(summary_text, summary_filename)
            summary_doc = self._upload_and_wait(
                file_path=summary_path,
                store_name=store_name,
                display_name=Path(summary_filename).stem,
                deadline=self._deadline_from_start(start),
            )
            summary_doc_name = summary_doc["name"]

        return {
            "store_doc_name": original_doc["name"],
            "summary_doc_name": summary_doc_name,
        }

    def cleanup_store(self, store_name: str) -> list[dict[str, str]]:
        documents = self.list_documents(store_name)
        for document in documents:
            self.delete_document(document["name"])
        logger.warning(
            "cleanup store executed store_id=%s deleted_count=%s",
            store_name,
            len(documents),
        )
        return documents

    def delete_document(self, doc_name: str) -> None:
        try:
            self.client.file_search_stores.documents.delete(
                name=doc_name,
                config={"force": True},
            )
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        logger.info("deleted document resource_name=%s", doc_name)

    def query_store(
        self,
        *,
        store_name: str,
        question: str,
        model: str = "gemini-2.5-flash",
    ) -> dict[str, Any]:
        try:
            response = self.client.models.generate_content(
                model=model,
                contents=question,
                config=types.GenerateContentConfig(
                    tools=[
                        types.Tool(
                            file_search=types.FileSearch(
                                file_search_store_names=[store_name]
                            )
                        )
                    ]
                ),
            )
        except genai_errors.APIError as exc:
            raise RagServiceError(str(exc)) from exc
        logger.info("query executed store_id=%s model=%s", store_name, model)
        return {
            "answer": response.text,
            "sources": self._parse_grounding_metadata(response),
        }

    def _upload_and_wait(
        self,
        *,
        file_path: str | Path,
        store_name: str,
        display_name: str,
        deadline: float,
    ) -> dict[str, str]:
        operation = self.client.file_search_stores.upload_to_file_search_store(
            file=str(file_path),
            file_search_store_name=store_name,
            config={"display_name": display_name},
        )

        while not operation.done:
            if time.perf_counter() >= deadline:
                raise RagUploadTimeoutError(
                    "Upload timed out. Google API did not respond within 90 seconds."
                )
            time.sleep(1)
            operation = self.client.operations.get(operation)

        documents = list(self.client.file_search_stores.documents.list(parent=store_name))
        for document in documents:
            if getattr(document, "display_name", "") == display_name:
                return {
                    "name": document.name,
                    "display_name": document.display_name,
                }

        raise RagServiceError(f"Uploaded document '{display_name}' not found in store.")

    def _deadline_from_start(self, start: float) -> float:
        return start + UPLOAD_TIMEOUT_SECONDS

    def _write_summary_file(self, summary_text: str, summary_filename: str) -> Path:
        summary_path = Path("uploads") / summary_filename
        summary_path.write_text(summary_text, encoding="utf-8")
        return summary_path

    def _parse_grounding_metadata(self, response: Any) -> list[dict[str, str]]:
        try:
            candidate = response.candidates[0]
        except (AttributeError, IndexError, TypeError):
            return []

        metadata = getattr(candidate, "grounding_metadata", None)
        if metadata is None:
            return []

        chunks = getattr(metadata, "grounding_chunks", None) or []
        supports = getattr(metadata, "grounding_supports", None) or []

        support_by_chunk: dict[int, list[str]] = {}
        for support in supports:
            indices = getattr(support, "grounding_chunk_indices", None) or []
            text = getattr(getattr(support, "segment", None), "text", "") or ""
            for index in indices:
                if text:
                    support_by_chunk.setdefault(index, []).append(text)

        sources: list[dict[str, str]] = []
        for index, chunk in enumerate(chunks):
            web_info = getattr(chunk, "web", None)
            retrieved_context = getattr(chunk, "retrieved_context", None)
            title = (
                getattr(retrieved_context, "title", None)
                or getattr(web_info, "title", None)
                or "Unknown source"
            )
            snippets = support_by_chunk.get(index, [])
            sources.append(
                {
                    "doc_name": title,
                    "chunk_text": snippets[0] if snippets else "",
                    "relevance_score": None,
                }
            )
        return sources

    def _to_jsonable(self, value: Any) -> Any:
        if value is None or isinstance(value, (str, int, float, bool)):
            return value
        if isinstance(value, list):
            return [self._to_jsonable(item) for item in value]
        if isinstance(value, dict):
            return {str(key): self._to_jsonable(item) for key, item in value.items()}
        if hasattr(value, "model_dump"):
            return self._to_jsonable(value.model_dump())
        if hasattr(value, "__dict__"):
            return {
                key: self._to_jsonable(item)
                for key, item in vars(value).items()
                if not key.startswith("_")
            }
        return str(value)


rag_service = RagService()
